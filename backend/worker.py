import os
import requests
import io
import json
import time
import hmac
import hashlib
from datetime import datetime, timezone
from celery import Celery
from pptx import Presentation as PPTXPresentation
import openai
from redis import Redis
from statuses import PresentationStatus

REDIS_URL = os.getenv("REDIS_URL", "redis://localhost:6379/0")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "")
INTERNAL_SERVICE_TOKEN = os.getenv("INTERNAL_SERVICE_TOKEN", "")
INTERNAL_SIGNING_SECRET = os.getenv("INTERNAL_SIGNING_SECRET", "")
REDIS_PUBSUB_CHANNEL = os.getenv("REDIS_PUBSUB_CHANNEL", "presentation_events")

if OPENAI_API_KEY:
    openai.api_key = OPENAI_API_KEY

celery_app = Celery("slideflow_worker", broker=REDIS_URL, backend=REDIS_URL)

celery_app.conf.update(
    task_serializer="json",
    accept_content=["json"],
    result_serializer="json",
    timezone="UTC",
    enable_utc=True,
)


def publish_event(event: str, presentation_id: str, owner_id: str, data: dict):
    if not presentation_id:
        return
    try:
        redis_client = Redis.from_url(REDIS_URL, decode_responses=True)
        payload = {
            "event": event,
            "version": 1,
            "timestamp": datetime.now(timezone.utc).isoformat(),
            "presentationId": presentation_id,
            "ownerId": owner_id,
            "data": data,
        }
        redis_client.publish(REDIS_PUBSUB_CHANNEL, json.dumps(payload))
        redis_client.close()
    except Exception as exc:
        print(f"[Worker] Failed to publish event {event}: {exc}")

def generate_ai_summary(text: str) -> str:
    if not OPENAI_API_KEY or not text.strip():
        return f"Summary: {text[:80]}..." if text.strip() else "No text found on this slide."
    try:
        response = openai.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are a professional presentation assistant. Summarize the following slide text in 1-2 concise sentences. Also extract 3-5 keywords. Format output as: Summary: [summary]\nKeywords: [kw1, kw2, ...]"},
                {"role": "user", "content": f"Slide text:\n{text}"}
            ],
            max_tokens=150
        )
        return response.choices[0].message.content
    except Exception as e:
        print("OpenAI Error:", e)
        return f"Summary: {text[:80]}..."

def build_internal_request(payload: dict) -> tuple[str, dict]:
    body = json.dumps(payload, separators=(",", ":"), ensure_ascii=False)
    headers = {
        "Content-Type": "application/json",
    }

    if INTERNAL_SERVICE_TOKEN:
        headers["x-service-token"] = INTERNAL_SERVICE_TOKEN

    if INTERNAL_SIGNING_SECRET:
        timestamp = str(int(time.time()))
        signed_payload = f"{timestamp}.{body}".encode("utf-8")
        signature = hmac.new(
            INTERNAL_SIGNING_SECRET.encode("utf-8"),
            signed_payload,
            hashlib.sha256,
        ).hexdigest()
        headers["x-service-timestamp"] = timestamp
        headers["x-service-signature"] = signature

    return body, headers


def update_presentation_status(next_api_url: str, presentation_id: str, status: str):
    if not presentation_id:
        return
    payload = {"presentationId": presentation_id, "status": status}
    body, headers = build_internal_request(payload)
    res = requests.patch(
        f"{next_api_url}/api/internal/presentations",
        data=body,
        headers=headers,
        timeout=10,
    )
    res.raise_for_status()

@celery_app.task(name="parse_presentation")
def parse_presentation_task(
    file_path: str,
    owner_id: str,
    file_name: str,
    next_api_url: str = "http://localhost:3000",
    presentation_id: str = "",
):
    print(f"[Worker] Starting parsing for {file_name} (presentation_id={presentation_id or 'new'})")
    is_local_file = not file_path.startswith("http")

    try:
        if presentation_id:
            publish_event(
                event="presentation.processing_started",
                presentation_id=presentation_id,
                owner_id=owner_id,
                data={"status": PresentationStatus.PROCESSING.value, "fileName": file_name, "totalSlides": None},
            )

        # 1. Load file bytes
        if is_local_file:
            with open(file_path, "rb") as f:
                file_bytes = io.BytesIO(f.read())
        else:
            response = requests.get(file_path, timeout=30)
            response.raise_for_status()
            file_bytes = io.BytesIO(response.content)

        # 2. Parse PPTX slides
        prs = PPTXPresentation(file_bytes)
        total_slides = len(prs.slides)
        slides_data = []
        for i, slide in enumerate(prs.slides):
            text_parts = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
            content = " ".join(text_parts).strip()
            ai_summary = generate_ai_summary(content)
            slides_data.append({
                "order": i + 1,
                "content": content,
                "aiSummary": ai_summary,
                "keywords": [],
            })
            if presentation_id:
                progress_percent = int(((i + 1) / max(total_slides, 1)) * 90)
                publish_event(
                    event="presentation.slide_parsed",
                    presentation_id=presentation_id,
                    owner_id=owner_id,
                    data={
                        "status": PresentationStatus.PROCESSING.value,
                        "slideOrder": i + 1,
                        "totalSlides": total_slides,
                        "progressPercent": progress_percent,
                        "slidePreview": {
                            "aiSummary": ai_summary[:160],
                            "keywords": [],
                        },
                    },
                )

        # 3. Create OR update the presentation record in Prisma via Next.js internal API
        if presentation_id:
            # Presentation record already created by Next.js upload-proxy — just update slides
            print(f"[Worker] Using existing presentation_id={presentation_id}")
        else:
            # No existing record — create one
            presentation_payload = {
                "title": file_name.rsplit(".", 1)[0],
                "fileUrl": file_path if not is_local_file else "",
                "ownerId": owner_id,
            }
            presentation_body, presentation_headers = build_internal_request(presentation_payload)
            pres_res = requests.post(
                f"{next_api_url}/api/internal/presentations",
                data=presentation_body,
                headers=presentation_headers,
                timeout=10,
            )
            pres_res.raise_for_status()
            presentation_id = pres_res.json()["id"]
            print(f"[Worker] Created new presentation_id={presentation_id}")
            publish_event(
                event="presentation.processing_started",
                presentation_id=presentation_id,
                owner_id=owner_id,
                data={"status": PresentationStatus.PROCESSING.value, "fileName": file_name, "totalSlides": total_slides},
            )

        # Attach presentation ID to each slide
        for s in slides_data:
            s["presentationId"] = presentation_id

        publish_event(
            event="presentation.progress",
            presentation_id=presentation_id,
            owner_id=owner_id,
            data={
                "status": PresentationStatus.PROCESSING.value,
                "stage": "persisting_slides",
                "processedSlides": len(slides_data),
                "totalSlides": total_slides,
                "progressPercent": 95,
                "message": "Saving parsed slides",
            },
        )

        # 4. Save slides + mark as PARSED
        slides_payload = {
            "slides": slides_data,
            "presentationId": presentation_id,
        }
        slides_body, slides_headers = build_internal_request(slides_payload)
        slides_res = requests.post(
            f"{next_api_url}/api/internal/slides",
            data=slides_body,
            headers=slides_headers,
            timeout=10,
        )
        slides_res.raise_for_status()

        publish_event(
            event="presentation.completed",
            presentation_id=presentation_id,
            owner_id=owner_id,
            data={
                "status": PresentationStatus.PARSED.value,
                "processedSlides": len(slides_data),
                "totalSlides": total_slides,
                "progressPercent": 100,
            },
        )
        print(f"[Worker] Done. Parsed {len(slides_data)} slides.")
        return {"status": "success", "presentation_id": presentation_id, "slides": len(slides_data)}

    except Exception as e:
        print(f"[Worker] Error: {str(e)}")
        try:
            if presentation_id:
                update_presentation_status(next_api_url, presentation_id, PresentationStatus.ERROR.value)
                publish_event(
                    event="presentation.failed",
                    presentation_id=presentation_id,
                    owner_id=owner_id,
                    data={
                        "status": PresentationStatus.ERROR.value,
                        "stage": "parsing",
                        "errorCode": "PRESENTATION_PARSE_ERROR",
                        "message": "Failed to parse presentation file.",
                        "retryable": False,
                    },
                )
        except Exception as status_exc:
            print(f"[Worker] Failed to update ERROR status: {status_exc}")
        return {"status": "error", "message": str(e)}
    finally:
        if is_local_file and os.path.exists(file_path):
            os.remove(file_path)
            print(f"[Worker] Cleaned up temp file: {file_path}")

